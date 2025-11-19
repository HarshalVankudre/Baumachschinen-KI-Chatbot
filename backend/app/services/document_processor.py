"""
Document Processing Service

Handles document processing workflow:
- OCR extraction with Docling
- Text chunking with semantic boundaries
- Embedding generation with OpenAI
- Vector storage in Weaviate (with multi-tenancy and hybrid search)
- MongoDB status updates
"""

import logging
import os
import time
import warnings
import gc
import asyncio
from typing import List, Dict, Any, Optional, Callable, Coroutine
from datetime import datetime
import uuid
from pathlib import Path

# Fix PIL DecompressionBombWarning for large PDF images
# PDFs converted to high-resolution images can exceed PIL's default 89MP limit
# Increase to 200MP to handle large technical diagrams/flowcharts
from PIL import Image

# Document processing libraries
try:
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        EasyOcrOptions,
        TableFormerMode,
        PictureDescriptionApiOptions
    )
    from docling.datamodel.accelerator_options import AcceleratorOptions, AcceleratorDevice
    from docling.datamodel.settings import settings as docling_settings
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    logging.warning("Docling not available. Document processing will be limited.")

# Text processing
import tiktoken

from app.config import get_settings
from app.core.database import get_database
from app.services.openai_service import get_openai_service
from app.services.weaviate_service import weaviate_service
from app.services.document_events import get_document_events_manager
from app.services.vision_extraction_service import get_vision_extraction_service
from app.services.text_chunker import get_text_chunker

logger = logging.getLogger(__name__)
settings = get_settings()

# Constants
DEFAULT_CHUNK_SIZE = 500
DEFAULT_CHUNK_OVERLAP = 50
BATCH_SIZE_EMBEDDINGS = 100
BATCH_SIZE_WEAVIATE = 1000  # Weaviate supports larger batches than Pinecone (1000 vs 100)
HEARTBEAT_INTERVAL_SECONDS = 15
MAX_METADATA_TEXT_LENGTH = 10000  # Weaviate supports longer metadata text
PIL_MAX_IMAGE_PIXELS = 200_000_000  # 200 megapixels


def _cleanup_memory(logger_func=None):
    """
    Aggressive memory cleanup after OCR processing

    Releases memory from EasyOCR/PyTorch models that can consume 18GB+ RAM.
    Should be called after document processing completes.
    """
    if logger_func:
        logger_func("Starting aggressive memory cleanup...")

    # Force Python garbage collection
    gc.collect()

    # Clear PyTorch CUDA cache if available
    try:
        import torch
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
            torch.cuda.synchronize()
            if logger_func:
                logger_func("Cleared PyTorch CUDA cache")
    except ImportError:
        pass  # PyTorch not installed
    except Exception as e:
        if logger_func:
            logger_func(f"Warning: Could not clear CUDA cache: {e}")

    # Additional garbage collection pass
    gc.collect()

    if logger_func:
        logger_func("Memory cleanup completed")


class DocumentProcessor:
    """
    Service for processing uploaded documents.

    Handles the complete document processing pipeline:
    1. OCR extraction with Docling
    2. Text chunking with semantic boundaries
    3. Embedding generation with OpenAI
    4. Vector storage in Weaviate (with multi-tenancy support)
    5. MongoDB status updates

    Attributes:
        openai_service: OpenAI service for embeddings
        weaviate_service: Weaviate service for vector storage
        db: MongoDB database connection
        events_manager: Document events manager for SSE
        encoding: Tiktoken encoding for token counting
        converter: Docling document converter (if available)
    """

    def __init__(self):
        """Initialize document processor with service dependencies and configure Docling."""
        # Set PIL image limit before any image processing
        Image.MAX_IMAGE_PIXELS = PIL_MAX_IMAGE_PIXELS

        # Initialize services
        self.openai_service = get_openai_service()
        self.weaviate_service = weaviate_service
        self.db = get_database()
        self.events_manager = get_document_events_manager()
        self.vision_extraction_service = get_vision_extraction_service()
        self.vision_temp_dir = Path("temp_uploads/vision")
        try:
            self.vision_temp_dir.mkdir(parents=True, exist_ok=True)
        except Exception as vision_dir_error:
            logger.warning("Falling back to /tmp for vision tiles: %s", vision_dir_error)
            self.vision_temp_dir = Path("/tmp")

        # Initialize text chunker
        self.text_chunker = get_text_chunker()

        # Initialize Docling converter with advanced optimization if available
        if DOCLING_AVAILABLE:
            # Suppress PyTorch pin_memory warning when no GPU is available
            warnings.filterwarnings("ignore", message=".*pin_memory.*no accelerator.*")

            # Configure CPU optimization - use all available cores
            import multiprocessing
            num_cores = multiprocessing.cpu_count()

            logger.info(f"Configuring Docling with {num_cores} CPU cores for maximum performance")

            # Configure accelerator options for maximum CPU utilization
            accelerator_options = AcceleratorOptions(
                num_threads=num_cores,  # Use all CPU cores
                device=AcceleratorDevice.CPU  # Force CPU (works on all systems)
            )

            # Configure advanced OCR options for better quality
            ocr_options = EasyOcrOptions(
                force_full_page_ocr=True,  # Full page OCR for maximum quality
                use_gpu=False,  # Use CPU (GPU requires specific setup)
            )

            # Configure GPT-4 Vision API for chart/diagram/flowchart extraction
            # Get OpenAI API key from settings
            openai_api_key = settings.openai_api_key

            if not openai_api_key:
                raise ValueError("OpenAI API key not configured in settings")

            gpt4_vision_options = PictureDescriptionApiOptions(
                url="https://api.openai.com/v1/chat/completions",
                params=dict(
                    model="gpt-4o",  # GPT-4o with vision capabilities
                    max_tokens=8192,  # Balanced token limit for detailed extraction without overwhelming the API
                ),
                headers={
                    "Authorization": f"Bearer {openai_api_key}"
                },
                prompt=(
                    "Extrahiere ALLE Informationen aus diesem Bild mit MAXIMALER Detailgenauigkeit:\n\n"

                    "1. TEXTEXTRAKTION (VOLLSTÄNDIG):\n"
                    "   - Transkribiere JEDES einzelne sichtbare Zeichen ohne Ausnahme\n"
                    "   - ALLE Überschriften, Unterüberschriften, Beschriftungen und Anmerkungen\n"
                    "   - ALLE Zahlen, Werte, Einheiten, Messungen und ihre Kontexte\n"
                    "   - ALLE Legenden, Schlüssel, Fußnoten und Randbemerkungen\n"
                    "   - ALLE Wasserzeichen, Logos und Markierungen mit Text\n"
                    "   - Selbst teilweise sichtbare oder schwer lesbare Texte (mit Hinweis auf Unsicherheit)\n\n"

                    "2. TECHNISCHE SPEZIFIKATIONEN (BAUMASCHINEN-FOKUS):\n"
                    "   - ALLE Maschinenbezeichnungen, Modellnummern, Typenbezeichnungen\n"
                    "   - ALLE technische Daten: Motorleistung (kW/PS), Gewicht, Abmessungen\n"
                    "   - ALLE Kapazitäten: Schaufelvolumen, Nutzlast, Hubkraft, Reichweite\n"
                    "   - ALLE Betriebsparameter: Geschwindigkeit, Steigfähigkeit, Wendekreis\n"
                    "   - ALLE Hydraulikdaten: Druck, Durchfluss, Zylindermaße\n"
                    "   - ALLE Herstellerinformationen, Seriennummern, Baujahre\n"
                    "   - ALLE Wartungsintervalle, Ölmengen, Füllmengen\n\n"

                    "3. DIAGRAMME & GRAPHEN (EXAKT):\n"
                    "   - Diagrammtyp präzise benennen\n"
                    "   - ALLE Achsenbeschriftungen mit vollständigen Einheiten\n"
                    "   - JEDEN EINZELNEN Datenpunkt mit exaktem Wert (nicht geschätzt)\n"
                    "   - ALLE Kurvenverläufe detailliert beschreiben\n"
                    "   - ALLE Legende-Einträge und Farbcodierungen\n"
                    "   - Trends, Maxima, Minima, Wendepunkte beschreiben\n\n"

                    "4. TABELLEN (STRUKTURIERT):\n"
                    "   - ALLE Spaltenüberschriften und Zeilenbeschriftungen vollständig\n"
                    "   - JEDEN EINZELNEN Zellenwert ohne Ausnahme\n"
                    "   - Tabellenstruktur klar mit Markdown oder strukturiertem Text\n"
                    "   - ALLE Summen, Zwischensummen, Durchschnitte, Berechnungen\n"
                    "   - Einheiten und Fußnoten zu Tabellenwerten\n\n"

                    "5. DIAGRAMME & FLUSSDIAGRAMME:\n"
                    "   - JEDE Form, Kästchen, Symbol mit vollständigem Textinhalt\n"
                    "   - ALLE Pfeile mit Richtung und Beschriftung\n"
                    "   - ALLE Verbindungen zwischen Komponenten\n"
                    "   - Prozessschritte nummeriert und in Reihenfolge\n"
                    "   - Entscheidungspunkte mit Bedingungen\n\n"

                    "6. BILDER & FOTOS:\n"
                    "   - Detaillierte Beschreibung der abgebildeten Maschine/Komponente\n"
                    "   - Sichtbare Teile, Baugruppen, Merkmale benennen\n"
                    "   - Zustand, Farbe, besondere Kennzeichen beschreiben\n\n"

                    "7. FORMELN & MATHEMATIK:\n"
                    "   - ALLE mathematischen Ausdrücke vollständig\n"
                    "   - Variablen und ihre Bedeutung erklären\n"
                    "   - Berechnungsbeispiele falls vorhanden\n\n"

                    "WICHTIG:\n"
                    "- Sei ausführlich und präzise\n"
                    "- Fasse NICHTS zusammen - extrahiere ALLES\n"
                    "- Bei Unsicherheit: schreibe den Text trotzdem mit Hinweis [unsicher]\n"
                    "- Ausgabe komplett auf Deutsch\n"
                    "- Strukturiere die Ausgabe klar und übersichtlich"
                ),
                scale=3.0,  # Maximum quality image processing (3x resolution)
                timeout=600,  # 10 minutes per image for very large/complex documents
                batch_size=2,  # Process 2 images at a time (GPT-4 API rate limit consideration)
            )

            # Configure PDF pipeline with advanced options
            pipeline_options = PdfPipelineOptions(
                accelerator_options=accelerator_options,
                do_ocr=True,  # Enable OCR for text (GPT-4 Vision also has OCR, but traditional OCR is faster for pure text)
                do_table_structure=True,  # Enable table extraction
                do_picture_description=True,  # Enable GPT-4 Vision for charts/diagrams
                ocr_options=ocr_options,
                picture_description_options=gpt4_vision_options,  # GPT-4 Vision configuration
                generate_page_images=False,  # Don't generate full page images (saves memory)
                generate_picture_images=True,  # Generate images for Vision AI
                images_scale=3.0,  # Maximum quality (3x scale) for GPT-4 Vision
            )

            # Enable remote services for API-based vision models
            pipeline_options.enable_remote_services = True  # Required for Vision AI

            # Configure table extraction for ACCURATE mode (better quality)
            pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE
            pipeline_options.table_structure_options.do_cell_matching = True

            # Enable pipeline profiling for performance monitoring
            docling_settings.debug.profile_pipeline_timings = True

            # Create converter with optimized configuration
            self.converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_options=pipeline_options
                    )
                }
            )

            logger.info(
                f"Docling DocumentConverter initialized with:\n"
                f"  - CPU cores: {num_cores}\n"
                f"  - Full page OCR: Enabled (EasyOCR)\n"
                f"  - Table extraction: ACCURATE mode\n"
                f"  - Image scale: 3.0x for maximum quality\n"
                f"  - Vision AI: GPT-4o (OpenAI API) for charts/diagrams/flowcharts\n"
                f"  - GPT-4 max tokens: 8192 (detailed extraction)\n"
                f"  - GPT-4 timeout: 600s (10 min per image)\n"
                f"  - PIL image limit: 200MP (handles large diagrams)\n"
                f"  - Remote services: ENABLED for API-based vision models\n"
                f"  - Comprehensive extraction: ALL text, tables, charts, diagrams, and technical content"
            )
        else:
            self.converter = None
            logger.warning("Docling not available - using fallback text extraction")

    async def process_document(
        self,
        document_id: str,
        file_path: str,
        category: str,
        uploader_name: str
    ) -> Dict[str, Any]:
        """
        Process document end-to-end

        Steps:
        1. Extract text with Docling OCR
        2. Chunk text into ~500 token segments
        3. Generate embeddings for each chunk
        4. Store vectors in Pinecone
        5. Update MongoDB with completion status

        Args:
            document_id: Unique document ID
            file_path: Path to uploaded file
            category: Document category
            uploader_name: Username of uploader

        Returns:
            Processing result with status and metadata
        """
        start_time = time.time()

        try:
            logger.info(f"Starting document processing for {document_id}")

            # Update status to processing
            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {"$set": {
                    "processing_status": "processing",
                    "processing_step": "extracting_text",
                    "processing_progress": 0
                }}
            )

            # Broadcast processing started
            await self.events_manager.broadcast_progress(
                document_id=document_id,
                status="processing",
                step="extracting_text",
                progress=0
            )

            # Step 1: Extract text from document with progress heartbeat
            logger.info(f"Extracting text from {file_path}")

            # Use progress-aware extraction for better UX during long OCR operations
            try:
                text_content = await self._extract_text_with_progress(file_path, document_id)
            except Exception as e:
                # Fall back to standard extraction with fallback methods
                logger.warning(f"Progress-aware extraction failed, using fallback: {str(e)}")
                text_content = await self._extract_text(file_path)

            if not text_content or len(text_content.strip()) < 10:
                raise ValueError("No text content extracted from document")

            logger.info(f"Extracted {len(text_content)} characters from document")

            # Update progress after extraction
            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {"$set": {
                    "processing_step": "chunking",
                    "processing_progress": 30
                }}
            )

            # Step 2: Chunk text into segments
            logger.info(f"[CHUNKING] Starting text chunking: {len(text_content)} characters")
            chunks = self.text_chunker.chunk_text(text_content, chunk_size=500, chunk_overlap=50)

            # Log chunk statistics
            total_chunk_chars = sum(len(chunk) for chunk in chunks)
            avg_chunk_size = total_chunk_chars // len(chunks) if chunks else 0
            logger.info(f"[CHUNKING] ✓ Created {len(chunks)} chunks (avg {avg_chunk_size} chars per chunk)")

            # Log first few chunks for debugging
            for i, chunk in enumerate(chunks[:3]):
                logger.info(f"[CHUNKING] Chunk #{i+1} preview: {chunk[:100]}...")

            # Update progress after chunking
            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {"$set": {
                    "processing_step": "generating_embeddings",
                    "processing_progress": 50
                }}
            )

            if not chunks:
                raise ValueError("No chunks created from document text")

            # Step 3: Generate embeddings for chunks (batch processing)
            logger.info(f"[EMBEDDINGS] Starting embedding generation for {len(chunks)} chunks")
            embeddings = await self._generate_embeddings_batch(chunks)
            embedding_dim = len(embeddings[0]) if embeddings and len(embeddings) > 0 else 0
            logger.info(f"[EMBEDDINGS] ✓ Generated {len(embeddings)} embeddings (dimension: {embedding_dim})")

            # Update progress after embeddings
            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {"$set": {
                    "processing_step": "storing_vectors",
                    "processing_progress": 80
                }}
            )

            # Step 4: Store vectors in Weaviate
            logger.info("Storing vectors in Weaviate")
            filename = os.path.basename(file_path)
            await self._store_vectors(
                document_id=document_id,
                filename=filename,
                category=category,
                uploader_name=uploader_name,
                chunks=chunks,
                embeddings=embeddings
            )

            # Step 5: Update MongoDB with success
            processing_time = time.time() - start_time
            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {
                    "$set": {
                        "processing_status": "completed",
                        "processing_step": None,  # Clear processing step
                        "processing_progress": 100,  # Set to 100% complete
                        "chunk_count": len(chunks),
                        "processing_time_seconds": round(processing_time, 2)
                    }
                }
            )

            logger.info(
                f"Document {document_id} processed successfully in {processing_time:.2f}s "
                f"({len(chunks)} chunks)"
            )

            # Broadcast completion
            await self.events_manager.broadcast_progress(
                document_id=document_id,
                status="completed",
                progress=100,
                chunk_count=len(chunks)
            )

            # Clean up temporary file
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    logger.info(f"Cleaned up temporary file: {file_path}")
            except Exception as e:
                logger.warning(f"Failed to clean up temporary file: {str(e)}")

            return {
                "status": "success",
                "document_id": document_id,
                "chunk_count": len(chunks),
                "processing_time_seconds": round(processing_time, 2)
            }

        except Exception as e:
            # Update MongoDB with failure
            processing_time = time.time() - start_time
            error_message = str(e)

            logger.error(f"Document processing failed for {document_id}: {error_message}")

            await self.db.document_metadata.update_one(
                {"document_id": document_id},
                {
                    "$set": {
                        "processing_status": "failed",
                        "error_message": error_message,
                        "processing_time_seconds": round(processing_time, 2)
                    }
                }
            )

            # Broadcast failure
            await self.events_manager.broadcast_progress(
                document_id=document_id,
                status="failed",
                error=error_message
            )

            # Clean up temporary file
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception:
                pass

            return {
                "status": "failed",
                "document_id": document_id,
                "error": error_message
            }

        finally:
            # Always cleanup memory after document processing (success or failure)
            # This releases EasyOCR/PyTorch models that can consume 18GB+ RAM
            logger.info("Running final memory cleanup after document processing")
            _cleanup_memory(lambda msg: logger.info(msg))

    async def _extract_text(self, file_path: str) -> str:
        """
        Extract text from document using Docling OCR

        Args:
            file_path: Path to document file

        Returns:
            Extracted text content
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()

        # Use Docling for comprehensive document conversion
        if DOCLING_AVAILABLE and self.converter:
            try:
                result = self.converter.convert(file_path)

                # Extract text content from conversion result
                text_parts = []

                # Docling provides structured document content
                # Extract text from document structure
                if hasattr(result, 'document'):
                    doc = result.document

                    # Get markdown representation (preserves structure)
                    if hasattr(doc, 'export_to_markdown'):
                        text_parts.append(doc.export_to_markdown())
                    # Fallback: get plain text
                    elif hasattr(doc, 'export_to_text'):
                        text_parts.append(doc.export_to_text())
                    # Last resort: convert to string
                    else:
                        text_parts.append(str(doc))

                text_content = "\n\n".join(text_parts)

                if text_content and len(text_content.strip()) > 0:
                    return text_content
                else:
                    logger.warning("Docling returned empty content, trying fallback extraction")

            except Exception as e:
                logger.error(f"Docling extraction failed: {str(e)}, trying fallback")

        # Fallback extraction methods
        return await self._fallback_extraction(file_path, file_ext)

    async def _extract_text_with_progress(
        self,
        file_path: str,
        document_id: str
    ) -> str:
        """
        Extract text with progress heartbeat updates for long-running OCR

        Runs Docling in asyncio thread pool to avoid blocking FastAPI's event loop.
        Uses asyncio.to_thread() which is simpler than multiprocessing and avoids
        pickling issues with complex class state.

        Sends heartbeat updates every 15 seconds to keep SSE connections alive
        during long OCR operations (5-10+ minutes for large PDFs).

        Args:
            file_path: Path to document file
            document_id: Document ID for progress broadcasting

        Returns:
            Extracted text content
        """
        import asyncio

        logger.info(f"[THREAD-BASED] Starting text extraction for {document_id} in background thread")

        # Use asyncio.to_thread() to run sync function in thread pool
        # This keeps event loop responsive for SSE without process pickling issues
        extraction_task = asyncio.create_task(
            asyncio.to_thread(self._extract_text_sync, file_path)
        )

        # Heartbeat loop: send progress updates while extraction runs
        progress_counter = 0
        heartbeat_interval = 15  # seconds

        while not extraction_task.done():
            try:
                # Wait with timeout to check periodically
                result = await asyncio.wait_for(
                    asyncio.shield(extraction_task),
                    timeout=heartbeat_interval
                )
                # If we get here, extraction completed
                logger.info(f"[THREAD-BASED] Completed text extraction for {document_id}: {len(result)} chars")
                return result

            except asyncio.TimeoutError:
                # Heartbeat: extraction still running
                progress_counter += heartbeat_interval
                logger.info(
                    f"[HEARTBEAT] Text extraction in progress for {document_id} "
                    f"({progress_counter}s elapsed)..."
                )

                # Update MongoDB with current progress (so frontend polling sees it)
                current_progress = min(15 + (progress_counter // 15) * 5, 85)
                await self.db.document_metadata.update_one(
                    {"document_id": document_id},
                    {"$set": {
                        "processing_step": "extracting_text",
                        "processing_progress": current_progress
                    }}
                )

                # Broadcast progress update via SSE
                await self.events_manager.broadcast_progress(
                    document_id=document_id,
                    status="processing",
                    step="extracting_text",
                    progress=current_progress
                )

        # Final result (if loop exits naturally)
        result = await extraction_task
        logger.info(f"[THREAD-BASED] Completed text extraction for {document_id}: {len(result)} chars")
        return result

    def _extract_text_sync(self, file_path: str) -> str:
        """
        Synchronous version of text extraction for thread pool execution

        This is needed because Docling's converter.convert() is synchronous
        and cannot be awaited directly. Runs in asyncio thread pool via
        asyncio.to_thread() to avoid blocking FastAPI's event loop.

        Args:
            file_path: Path to document file

        Returns:
            Extracted text content
        """
        import sys

        # Force flush logging to ensure visibility
        def log_and_flush(message):
            logger.info(message)
            print(f"[PROCESSOR] {message}", flush=True)
            sys.stdout.flush()

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path)

        log_and_flush(f"=" * 80)
        log_and_flush(f"STARTING TEXT EXTRACTION")
        log_and_flush(f"File: {os.path.basename(file_path)}")
        log_and_flush(f"Size: {file_size / 1024:.1f} KB")
        log_and_flush(f"Extension: {file_ext}")
        log_and_flush(f"=" * 80)

        # Use Docling for comprehensive document conversion
        if DOCLING_AVAILABLE and self.converter:
            try:
                log_and_flush("Step 1: Calling Docling converter.convert()...")
                log_and_flush("This may take several minutes for large PDFs with images...")

                # Time the conversion for performance monitoring
                conversion_start = time.time()

                log_and_flush("DOCLING CONVERSION STARTED - waiting for result...")
                result = self.converter.convert(file_path)
                conversion_time = time.time() - conversion_start

                # Log performance metrics
                log_and_flush(f"✓ DOCLING CONVERSION COMPLETED in {conversion_time:.2f}s")

                # Log timing details if available
                if hasattr(result, 'timings') and result.timings:
                    pipeline_times = result.timings.get("pipeline_total", {})
                    if hasattr(pipeline_times, 'times'):
                        logger.info(f"  Pipeline timings: {pipeline_times.times}")

                # DEBUG: Log document structure to understand what was processed
                if hasattr(result, 'document'):
                    doc = result.document
                    picture_count = 0
                    try:
                        from docling_core.types.doc import PictureItem
                        for element, level in doc.iterate_items():
                            if isinstance(element, PictureItem):
                                picture_count += 1
                        log_and_flush(f"Step 2: Found {picture_count} picture elements in document")
                    except Exception as e:
                        log_and_flush(f"WARNING: Could not count pictures: {str(e)}")

                # Extract text content from conversion result
                text_parts = []

                if hasattr(result, 'document'):
                    doc = result.document

                    # Export main document text (markdown format)
                    if hasattr(doc, 'export_to_markdown'):
                        markdown_text = doc.export_to_markdown()
                        text_parts.append(markdown_text)
                    elif hasattr(doc, 'export_to_text'):
                        text_parts.append(doc.export_to_text())
                    else:
                        text_parts.append(str(doc))

                    # CRITICAL FIX: Manually extract picture descriptions from GPT-4 Vision
                    # The export_to_markdown() doesn't include picture descriptions by default
                    log_and_flush("Step 3: Extracting GPT-4 Vision descriptions from pictures...")
                    picture_descriptions = []
                    try:
                        from docling_core.types.doc import PictureItem

                        # Iterate through all document items to find pictures
                        picture_item_count = 0
                        for element, level in doc.iterate_items():
                            if isinstance(element, PictureItem):
                                picture_item_count += 1
                                log_and_flush(f"  Processing PictureItem #{picture_item_count}: {element.self_ref if hasattr(element, 'self_ref') else 'unknown'}")
                                collected_chunks: List[str] = []

                                # Check if this picture has GPT-4 Vision descriptions
                                if hasattr(element, 'annotations') and element.annotations:
                                    log_and_flush(f"    ✓ Found {len(element.annotations)} annotation(s) - GPT-4 Vision WAS INVOKED!")
                                    for idx, annotation in enumerate(element.annotations):
                                        # DEBUG: Log annotation type and attributes
                                        annotation_type = type(annotation).__name__
                                        log_and_flush(f"    Annotation {idx}: type={annotation_type}")

                                        # Extract text from picture description annotation
                                        if hasattr(annotation, 'text') and annotation.text:
                                            chunk = f"\n[IMAGE DESCRIPTION - GPT-4 Vision]:\n{annotation.text}\n"
                                            collected_chunks.append(chunk)
                                            log_and_flush(f"    ✓ Extracted GPT-4 Vision text ({len(annotation.text)} chars)")
                                        # Alternative: check for 'description' attribute
                                        elif hasattr(annotation, 'description') and annotation.description:
                                            chunk = f"\n[IMAGE DESCRIPTION - GPT-4 Vision]:\n{annotation.description}\n"
                                            collected_chunks.append(chunk)
                                            log_and_flush(f"    ✓ Extracted GPT-4 Vision description ({len(annotation.description)} chars)")
                                        else:
                                            # Try to extract any text-like content
                                            log_and_flush(f"    WARNING: Annotation has no 'text' or 'description' attribute!")
                                            annotation_str = str(annotation)
                                            if annotation_str and len(annotation_str) > 10:
                                                chunk = f"\n[IMAGE DESCRIPTION - GPT-4 Vision]:\n{annotation_str}\n"
                                                collected_chunks.append(chunk)
                                                log_and_flush(f"    ✓ Extracted from str() ({len(annotation_str)} chars)")
                                else:
                                    log_and_flush(f"    ✗ PictureItem has NO annotations - GPT-4 Vision was NOT invoked!")

                                if collected_chunks:
                                    combined_description = "\n".join(collected_chunks)
                                    combined_description = self._maybe_enrich_picture_description(
                                        picture_item=element,
                                        base_text=combined_description,
                                        logger_func=log_and_flush
                                    )
                                    picture_descriptions.append(combined_description)
                                else:
                                    enriched_only = self._maybe_enrich_picture_description(
                                        picture_item=element,
                                        base_text="",
                                        logger_func=log_and_flush,
                                        force=True
                                    )
                                    if enriched_only:
                                        picture_descriptions.append(enriched_only)
                    except Exception as e:
                        log_and_flush(f"ERROR extracting picture descriptions: {str(e)}")
                        import traceback
                        log_and_flush(f"Traceback: {traceback.format_exc()}")

                    # Append all picture descriptions to the text content
                    if picture_descriptions:
                        text_parts.extend(picture_descriptions)
                        log_and_flush(f"✓ Step 4: Added {len(picture_descriptions)} GPT-4 Vision descriptions to text")
                    else:
                        log_and_flush("⚠ WARNING: No GPT-4 Vision descriptions found!")

                text_content = "\n\n".join(text_parts)

                if text_content and len(text_content.strip()) > 0:
                    log_and_flush(f"=" * 80)
                    log_and_flush(f"✓ EXTRACTION COMPLETE!")
                    log_and_flush(f"Total characters extracted: {len(text_content)}")
                    log_and_flush(f"Picture descriptions included: {len(picture_descriptions)}")
                    log_and_flush(f"=" * 80)

                    # Cleanup memory before returning (release EasyOCR/PyTorch ~18GB)
                    _cleanup_memory(log_and_flush)

                    return text_content
                else:
                    log_and_flush("ERROR: Docling returned empty content!")

            except Exception as e:
                log_and_flush(f"=" * 80)
                log_and_flush(f"✗ DOCLING EXTRACTION FAILED!")
                log_and_flush(f"Error: {str(e)}")
                import traceback
                log_and_flush(f"Traceback:\n{traceback.format_exc()}")
                log_and_flush(f"=" * 80)

                # Cleanup memory even on failure
                _cleanup_memory(log_and_flush)

        # If Docling fails or unavailable, raise error for async fallback
        # Cleanup memory before raising
        _cleanup_memory(lambda msg: logger.info(msg))
        raise ValueError("Docling extraction failed, fallback needed")

    async def _fallback_extraction(self, file_path: str, file_ext: str) -> str:
        """
        Fallback text extraction for when Docling is unavailable

        Args:
            file_path: Path to document
            file_ext: File extension

        Returns:
            Extracted text
        """
        text_content = ""

        # PDF extraction
        if file_ext == '.pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text_parts = []
                    for page in pdf_reader.pages:
                        text_parts.append(page.extract_text())
                    text_content = "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"PDF extraction failed: {str(e)}")

        # DOCX extraction
        elif file_ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                text_parts = [paragraph.text for paragraph in doc.paragraphs]
                text_content = "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"DOCX extraction failed: {str(e)}")

        # PPTX extraction
        elif file_ext in ['.ppt', '.pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text_content = "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"PPTX extraction failed: {str(e)}")

        # Excel extraction
        elif file_ext in ['.xls', '.xlsx']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True, data_only=True)
                text_parts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text_parts.append(row_text)
                text_content = "\n".join(text_parts)
            except Exception as e:
                logger.error(f"Excel extraction failed: {str(e)}")

        # Image files - note: basic OCR would require additional libraries
        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif']:
            logger.warning(f"Image file {file_ext} - OCR not implemented in fallback")
            text_content = f"[Image file: {os.path.basename(file_path)}]"

        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

        return text_content

    def _maybe_enrich_picture_description(
        self,
        picture_item,
        base_text: str,
        logger_func,
        force: bool = False
    ) -> str:
        """
        Use the tiling GPT-4 Vision service to extend picture descriptions.
        """
        service = getattr(self, "vision_extraction_service", None)
        if not service:
            return base_text

        current_length = len(base_text.strip()) if base_text else 0
        if not force and base_text and current_length >= service.min_characters:
            return base_text

        image_path = self._resolve_picture_image_path(picture_item)
        if not image_path:
            return base_text

        def _extract() -> str:
            return self._await_sync(lambda: service.extract_detailed_description(image_path))

        try:
            enriched = _extract()
        except Exception as exc:
            logger_func(f"    ⚠ Vision enrichment failed for {image_path}: {exc}")
            return base_text

        if not enriched:
            return base_text

        if enriched and len(enriched) > current_length:
            logger_func(
                f"    ⤴ Vision enrichment added {len(enriched)} chars "
                f"(prev {current_length}) from {os.path.basename(image_path)}"
            )
            if base_text:
                return (
                    f"{base_text}\n\n[GPT-4 Vision Enrichment | {len(enriched)} Zeichen]\n{enriched}"
                )
            return f"[GPT-4 Vision Enrichment | {len(enriched)} Zeichen]\n{enriched}"

        return base_text

    def _await_sync(self, coro_factory: Callable[[], Coroutine[Any, Any, Any]]) -> Any:
        """
        Run an async coroutine from a synchronous context (thread worker).
        """
        try:
            return asyncio.run(coro_factory())
        except RuntimeError:
            loop = asyncio.new_event_loop()
            try:
                asyncio.set_event_loop(loop)
                return loop.run_until_complete(coro_factory())
            finally:
                loop.run_until_complete(loop.shutdown_asyncgens())
                asyncio.set_event_loop(None)
                loop.close()

    def _resolve_picture_image_path(self, picture_item) -> Optional[str]:
        """
        Try to locate the image associated with a Docling PictureItem.
        """
        candidate_attrs = ["image_path", "image_filepath", "path", "file_path"]

        for attr in candidate_attrs:
            value = getattr(picture_item, attr, None)
            if isinstance(value, str) and os.path.exists(value):
                return value

        image_obj = getattr(picture_item, "image", None)
        if image_obj:
            for attr in candidate_attrs:
                value = getattr(image_obj, attr, None)
                if isinstance(value, str) and os.path.exists(value):
                    return value

            data = getattr(image_obj, "data", None)
            if isinstance(data, (bytes, bytearray)):
                temp_path = self.vision_temp_dir / f"picture_{uuid.uuid4().hex}.png"
                temp_path.write_bytes(data)
                return str(temp_path)
            if hasattr(data, "read"):
                temp_path = self.vision_temp_dir / f"picture_{uuid.uuid4().hex}.png"
                temp_path.write_bytes(data.read())
                return str(temp_path)

        return None

    async def _generate_embeddings_batch(
        self,
        chunks: List[str],
        batch_size: int = BATCH_SIZE_EMBEDDINGS
    ) -> List[List[float]]:
        """
        Generate embeddings for chunks in batches.

        Processes text chunks in batches to avoid API rate limits and
        memory issues with large documents.

        Args:
            chunks: List of text chunks
            batch_size: Batch size for API calls (default from BATCH_SIZE_EMBEDDINGS)

        Returns:
            List of embedding vectors (one per chunk)

        Raises:
            Exception: If embedding generation fails for any batch
        """
        all_embeddings = []
        total_batches = (len(chunks) + batch_size - 1) // batch_size

        # Process in batches
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            batch_num = i // batch_size + 1

            logger.info(
                f"[EMBEDDINGS] Processing batch {batch_num}/{total_batches} "
                f"({len(batch)} chunks)..."
            )

            try:
                embeddings = await self.openai_service.generate_embeddings_batch(batch)
                all_embeddings.extend(embeddings)

                logger.info(
                    f"[EMBEDDINGS] Batch {batch_num}/{total_batches} completed "
                    f"({len(embeddings)} embeddings generated)"
                )

            except Exception as e:
                logger.error(
                    f"Batch embedding generation failed for batch {batch_num}: {str(e)}",
                    exc_info=True
                )
                raise

        return all_embeddings

    async def _store_vectors(
        self,
        document_id: str,
        filename: str,
        category: str,
        uploader_name: str,
        chunks: List[str],
        embeddings: List[List[float]]
    ) -> None:
        """
        Store vectors in Weaviate with metadata.

        Prepares vector records with metadata and uploads them to Weaviate
        in batches for optimal performance.

        Args:
            document_id: Unique document identifier
            filename: Original filename
            category: Document category for filtering
            uploader_name: Username of uploader
            chunks: Text chunks corresponding to embeddings
            embeddings: Embedding vectors from OpenAI

        Raises:
            ValueError: If chunk and embedding counts don't match
            Exception: If Weaviate upload fails for any batch
        """
        if len(chunks) != len(embeddings):
            raise ValueError(
                f"Chunk count ({len(chunks)}) doesn't match embedding count "
                f"({len(embeddings)})"
            )

        # Prepare objects for Weaviate
        objects = self._prepare_weaviate_objects(
            document_id,
            filename,
            category,
            uploader_name,
            chunks,
            embeddings
        )

        # Upload to Weaviate in batches (use shared tenant for all users)
        await self._upload_objects_to_weaviate(objects, "shared")

    def _prepare_weaviate_objects(
        self,
        document_id: str,
        filename: str,
        category: str,
        uploader_name: str,
        chunks: List[str],
        embeddings: List[List[float]]
    ) -> List[Dict[str, Any]]:
        """
        Prepare object records for Weaviate upload.

        Args:
            document_id: Document identifier
            filename: Original filename
            category: Document category
            uploader_name: Uploader username
            chunks: Text chunks
            embeddings: Embedding vectors

        Returns:
            List of Weaviate object dictionaries
        """
        from datetime import datetime, timezone

        objects = []
        for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            # Weaviate auto-generates UUIDs, but we can specify them
            # Use deterministic UUID based on document_id and chunk_index
            import uuid
            object_uuid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"{document_id}_chunk{idx}"))

            # CRITICAL FIX: Use RFC3339 format for Weaviate DATE type
            # Weaviate requires proper timezone-aware datetime format
            timestamp = datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%S.%fZ')

            objects.append({
                "id": object_uuid,
                "vector": embedding,  # Pre-computed embedding
                "properties": {
                    "document_id": document_id,
                    "filename": filename,
                    "category": category,
                    "uploader_name": uploader_name,
                    "chunk_index": idx,
                    "text_content": chunk[:MAX_METADATA_TEXT_LENGTH],
                    "image_url": None,  # Optional: for multi-modal support
                    "created_at": timestamp
                }
            })

        return objects

    async def _upload_objects_to_weaviate(
        self,
        objects: List[Dict[str, Any]],
        tenant: Optional[str] = None
    ) -> None:
        """
        Upload objects to Weaviate in batches with multi-tenancy support.

        Args:
            objects: List of object dictionaries to upload
            tenant: Optional tenant name for multi-tenancy (uses uploader_name)

        Raises:
            Exception: If any batch upload fails
        """
        batch_size = BATCH_SIZE_WEAVIATE
        total_batches = (len(objects) + batch_size - 1) // batch_size

        logger.info(
            f"[WEAVIATE] Starting upload of {len(objects)} objects "
            f"in {total_batches} batches (batch_size={batch_size})..."
        )

        # Upload in batches (tenant creation handled in batch_upsert)
        result = await self.weaviate_service.batch_upsert(
            objects=objects,
            collection_name="Documents",
            tenant=tenant if settings.enable_weaviate_multitenancy else None,
            batch_size=batch_size
        )

        if result["failed"] > 0:
            logger.warning(
                f"[WEAVIATE] Upload completed with {result['failed']} failures "
                f"({result['success']} succeeded)"
            )
            # Log first few errors for debugging
            if "errors" in result and result["errors"]:
                logger.error(f"[WEAVIATE] Sample errors: {result['errors'][:3]}")

            # Raise exception if all objects failed
            if result["success"] == 0:
                raise Exception(
                    f"All {len(objects)} objects failed to upload. "
                    f"Errors: {result.get('errors', ['Unknown error'])[:3]}"
                )
        else:
            logger.info(
                f"[WEAVIATE] All {len(objects)} objects uploaded successfully!"
            )


# Singleton instances
_document_processor = None
_sycamore_processor = None


def get_document_processor() -> DocumentProcessor:
    """
    Get document processor instance based on configuration.

    Returns either the traditional Docling processor or the modern Aryn processor
    based on the USE_SYCAMORE_PROCESSOR setting.

    Feature Flag Control:
    - USE_SYCAMORE_PROCESSOR=true: Use Aryn SDK directly (recommended for better accuracy, no Sycamore)
    - USE_SYCAMORE_PROCESSOR=false: Use Docling + EasyOCR (legacy fallback)

    Returns:
        Document processor instance (either DocumentProcessor or ArynDocumentProcessor)
    """
    global _document_processor, _sycamore_processor

    # Check if Aryn processor should be used (replacing Sycamore)
    if settings.use_sycamore_processor:
        try:
            # Use Aryn SDK directly instead of Sycamore
            from app.services.aryn_processor import get_aryn_processor

            if _sycamore_processor is None:
                logger.info("[PROCESSOR] Initializing Aryn processor (Direct SDK)")
                _sycamore_processor = get_aryn_processor()

            return _sycamore_processor

        except ImportError as e:
            logger.warning(
                f"[PROCESSOR] Aryn processor not available ({str(e)}), "
                f"falling back to Docling processor. "
                f"Install with: pip install aryn-sdk"
            )
            # Fall through to Docling processor
        except Exception as e:
            logger.error(
                f"[PROCESSOR] Failed to initialize Aryn processor: {str(e)}, "
                f"falling back to Docling processor"
            )
            # Fall through to Docling processor

    # Use traditional Docling processor (fallback or explicit choice)
    if _document_processor is None:
        logger.info("[PROCESSOR] Initializing Docling processor (legacy)")
        _document_processor = DocumentProcessor()

    return _document_processor
